from __future__ import annotations

import hashlib
import io
from pathlib import Path
from typing import Any

from ..markdown import normalize_markdown
from ..models import KnowledgeDocument
from .base import DocumentProcessor


class PresentationProcessor(DocumentProcessor):
    """Processor for presentations (PPTX, ODP) converting slides and presenter notes into Markdown."""

    name = "presentation"
    supported_extensions = (".pptx", ".ppt", ".odp")
    supported_mimetypes = (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
        "application/vnd.oasis.opendocument.presentation",
    )

    def process_file(
        self,
        file_path: Path | str,
        *,
        metadata: dict[str, Any] | None = None,
        options: dict[str, Any] | None = None,
    ) -> KnowledgeDocument:
        path = Path(file_path)
        content_bytes = path.read_bytes()
        doc_metadata = dict(metadata or {})
        doc_metadata.setdefault("file_path", str(path.resolve()))
        doc_metadata.setdefault("filename", path.name)
        return self.process_bytes(
            content_bytes,
            filename=path.name,
            mime_type="",
            metadata=doc_metadata,
            options=options,
        )

    def process_bytes(
        self,
        content_bytes: bytes,
        *,
        filename: str = "",
        mime_type: str = "",
        metadata: dict[str, Any] | None = None,
        options: dict[str, Any] | None = None,
    ) -> KnowledgeDocument:
        del mime_type
        doc_metadata = dict(metadata or {})
        doc_options = dict(options or {})
        doc_id = hashlib.sha256(content_bytes).hexdigest()
        title = doc_metadata.get("title") or Path(filename).stem or "Apresentação"

        sections: list[str] = []
        slide_count = 0

        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content_bytes))
            slide_count = len(prs.slides)

            for index, slide in enumerate(prs.slides, start=1):
                slide_title = ""
                slide_lines: list[str] = []

                if slide.shapes.title and slide.shapes.title.text.strip():
                    slide_title = slide.shapes.title.text.strip()

                for shape in slide.shapes:
                    if shape == slide.shapes.title:
                        continue

                    # Text boxes and placeholders
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            p_text = paragraph.text.strip()
                            if not p_text:
                                continue
                            if paragraph.level > 0:
                                indent = "  " * paragraph.level
                                slide_lines.append(f"{indent}- {p_text}")
                            else:
                                slide_lines.append(p_text)

                    # Tables
                    if shape.has_table:
                        table = shape.table
                        rows = []
                        for row in table.rows:
                            row_cells = [cell.text.strip().replace("\n", " ").replace("|", "\\|") for cell in row.cells]
                            rows.append(row_cells)
                        if rows:
                            col_count = len(rows[0])
                            table_md = [
                                "| " + " | ".join(rows[0]) + " |",
                                "| " + " | ".join(["---"] * col_count) + " |",
                            ]
                            for row in rows[1:]:
                                padded = row + [""] * (col_count - len(row))
                                table_md.append("| " + " | ".join(padded[:col_count]) + " |")
                            slide_lines.append("\n".join(table_md))

                # Presenter Notes
                notes_text = ""
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()

                header = f"## Slide {index}"
                if slide_title:
                    header = f"## Slide {index} — {slide_title}"

                slide_body = "\n\n".join(slide_lines).strip()
                slide_section = f"{header}\n\n{slide_body}" if slide_body else header

                if notes_text and doc_options.get("include_notes", True):
                    slide_section += f"\n\n### Notas do apresentador\n\n{notes_text}"

                sections.append(slide_section)

        except Exception as exc:
            sections.append(f"*[Não foi possível processar apresentação: {exc}]*")

        full_content = normalize_markdown("\n\n".join(sections))
        doc_metadata.setdefault("raw_type", "presentation")
        doc_metadata.setdefault("slide_count", slide_count)
        doc_metadata.setdefault("filename", filename)

        return KnowledgeDocument(
            id=doc_id,
            container_id=str(doc_metadata.get("container_id") or "presentations"),
            title=title,
            content=full_content,
            original_url=str(doc_metadata.get("original_url") or doc_metadata.get("file_path") or filename),
            source_type=str(doc_metadata.get("source_type") or "local_files"),
            container_name=str(doc_metadata.get("container_name") or "Apresentações"),
            path=list(doc_metadata.get("path") or [title]),
            metadata=doc_metadata,
        )
